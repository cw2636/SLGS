#!/usr/bin/env python3
"""Generate SLGS EduMeet technical presentation for interviews."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colors ──
BG_DARK   = RGBColor(0x0A, 0x1A, 0x0E)
BG_CARD   = RGBColor(0x12, 0x2A, 0x18)
GOLD      = RGBColor(0xC9, 0xA2, 0x27)
GREEN     = RGBColor(0x22, 0xC5, 0x5E)
WHITE     = RGBColor(0xE0, 0xE0, 0xE0)
MUTED     = RGBColor(0x88, 0x88, 0x88)
RED       = RGBColor(0xEF, 0x44, 0x44)
BLUE      = RGBColor(0x3B, 0x82, 0xF6)

def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf

def add_bullet(tf, text, size=16, color=WHITE, level=0, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.level = level
    p.space_before = Pt(4)
    return p

def add_card(slide, left, top, width, height, color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

# ═══════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_bg(slide)

# Gold accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = GOLD; bar.line.fill.background()

add_text(slide, 1, 1.5, 11, 1.2, "SLGS EduMeet", size=48, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1, 2.7, 11, 0.8, "Virtual Classroom Platform for Sierra Leone Grammar School", size=24, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 1, 3.8, 11, 0.6, "React  •  Node.js/Express  •  Go (Gin)  •  LiveKit SFU  •  MongoDB  •  WebSocket", size=16, color=MUTED, align=PP_ALIGN.CENTER)

add_text(slide, 1, 5.5, 11, 0.5, "Christopher Wilson  |  April 2026", size=14, color=MUTED, align=PP_ALIGN.CENTER)

# Bottom bar
bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.42), prs.slide_width, Inches(0.08))
bar2.fill.solid(); bar2.fill.fore_color.rgb = GOLD; bar2.line.fill.background()


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: The Problem
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0.8, 0.4, 11, 0.8, "The Problem: Peer-to-Peer WebRTC Doesn't Scale", size=32, color=GOLD, bold=True)

# Left card: mesh diagram
add_card(slide, 0.8, 1.5, 5.5, 5.2)
tf = add_text(slide, 1.1, 1.7, 5, 0.5, "Mesh Topology (Original)", size=20, color=GOLD, bold=True)
tf = add_text(slide, 1.1, 2.3, 5, 4, "", size=15, color=WHITE)
add_bullet(tf, "Every participant connects to every other participant", size=15, color=WHITE)
add_bullet(tf, "O(n²) connections — 4 users = 6 connections", size=15, color=WHITE)
add_bullet(tf, "O(n) bandwidth — upload your stream N-1 times", size=15, color=WHITE)
add_bullet(tf, "8 users = 28 connections, 7× upload bandwidth", size=15, color=WHITE)
add_bullet(tf, "30 students = 435 connections (impossible)", size=15, color=RED, bold=True)
add_bullet(tf, "", size=10)
add_bullet(tf, "Real classroom failures:", size=16, color=GOLD, bold=True)
add_bullet(tf, "• CPU spikes from encoding every peer's stream", size=14, color=MUTED)
add_bullet(tf, "• Video freezes when 5th student joins", size=14, color=MUTED)
add_bullet(tf, "• Connection cascade failures", size=14, color=MUTED)
add_bullet(tf, "• No server-side recording possible", size=14, color=MUTED)

# Right card: connection table
add_card(slide, 6.8, 1.5, 5.7, 5.2)
tf = add_text(slide, 7.1, 1.7, 5.2, 0.5, "Connection Growth", size=20, color=GOLD, bold=True)
tf = add_text(slide, 7.1, 2.3, 5.2, 3, "", size=15, color=WHITE)
add_bullet(tf, "Participants    Connections    Upload/User", size=14, color=MUTED, bold=True)
add_bullet(tf, "     2                1             1 stream", size=14, color=GREEN)
add_bullet(tf, "     4                6             3 streams", size=14, color=WHITE)
add_bullet(tf, "     8              28             7 streams", size=14, color=WHITE)
add_bullet(tf, "    16             120            15 streams", size=14, color=RED)
add_bullet(tf, "    30             435            29 streams", size=14, color=RED, bold=True)
add_bullet(tf, "", size=10)
add_bullet(tf, "Formula: connections = n(n-1)/2", size=14, color=MUTED)
add_bullet(tf, "", size=10)
add_bullet(tf, "A Sierra Leone classroom has 30-50 students.", size=15, color=WHITE, bold=True)
add_bullet(tf, "Mesh topology is fundamentally unsuitable.", size=15, color=RED, bold=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: The Solution — LiveKit SFU
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0.8, 0.4, 11, 0.8, "The Solution: LiveKit SFU Architecture", size=32, color=GOLD, bold=True)

# Left: SFU explanation
add_card(slide, 0.8, 1.5, 5.5, 5.2)
tf = add_text(slide, 1.1, 1.7, 5, 0.5, "Selective Forwarding Unit", size=20, color=GREEN, bold=True)
tf = add_text(slide, 1.1, 2.3, 5, 4, "", size=15, color=WHITE)
add_bullet(tf, "Each user uploads ONE stream to the server", size=15, color=WHITE, bold=True)
add_bullet(tf, "Server forwards to all subscribers (no re-encoding)", size=15, color=WHITE)
add_bullet(tf, "", size=8)
add_bullet(tf, "Participants    Upload       Download", size=14, color=MUTED, bold=True)
add_bullet(tf, "     2           1 stream     1 stream", size=14, color=GREEN)
add_bullet(tf, "     8           1 stream     7 streams", size=14, color=GREEN)
add_bullet(tf, "    30           1 stream     adaptive", size=14, color=GREEN, bold=True)
add_bullet(tf, "   100           1 stream     adaptive", size=14, color=GREEN, bold=True)
add_bullet(tf, "", size=8)
add_bullet(tf, "Key: upload is always 1 stream regardless of room size", size=14, color=GOLD)

# Right: Why LiveKit
add_card(slide, 6.8, 1.5, 5.7, 5.2)
tf = add_text(slide, 7.1, 1.7, 5.2, 0.5, "Why LiveKit Specifically", size=20, color=GREEN, bold=True)
tf = add_text(slide, 7.1, 2.3, 5.2, 4, "", size=15, color=WHITE)
add_bullet(tf, "✓ Single Go binary — no Java/Docker complexity", size=15, color=GREEN)
add_bullet(tf, "✓ One YAML config file", size=15, color=GREEN)
add_bullet(tf, "✓ First-class React SDK (@livekit/components-react)", size=15, color=GREEN)
add_bullet(tf, "✓ Built-in TURN server (NAT traversal)", size=15, color=GREEN)
add_bullet(tf, "✓ Simulcast — server chooses quality per subscriber", size=15, color=GREEN)
add_bullet(tf, "✓ Server-side recording capability", size=15, color=GREEN)
add_bullet(tf, "✓ JWT token authentication (HMAC-SHA256)", size=15, color=GREEN)
add_bullet(tf, "", size=8)
add_bullet(tf, "vs Jitsi: complex stack (Java + XMPP + Oracleserver)", size=13, color=MUTED)
add_bullet(tf, "vs Mediasoup: library only, build everything yourself", size=13, color=MUTED)
add_bullet(tf, "vs Zoom SDK: proprietary, per-minute pricing", size=13, color=MUTED)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: System Architecture
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0.8, 0.4, 11, 0.8, "System Architecture — Three Coordinated Services", size=32, color=GOLD, bold=True)

# Browser box
add_card(slide, 1.5, 1.4, 10.3, 1.3, RGBColor(0x1A, 0x2A, 0x3A))
add_text(slide, 1.8, 1.5, 9.7, 0.4, "Browser (React)", size=18, color=BLUE, bold=True)
add_text(slide, 1.8, 1.9, 3, 0.4, "EduMeetLive.js", size=13, color=WHITE)
add_text(slide, 4.8, 1.9, 3, 0.4, "LiveKit Room (Video/Audio)", size=13, color=WHITE)
add_text(slide, 8.2, 1.9, 3, 0.4, "WebSocket Hook (Features)", size=13, color=WHITE)

# Three server boxes
add_card(slide, 1.5, 3.2, 3, 2.8)
tf = add_text(slide, 1.7, 3.3, 2.6, 0.4, "Express.js :3000", size=16, color=GOLD, bold=True)
tf = add_text(slide, 1.7, 3.8, 2.6, 2, "", size=13, color=WHITE)
add_bullet(tf, "Authentication", size=13)
add_bullet(tf, "Session management", size=13)
add_bullet(tf, "MongoDB queries", size=13)
add_bullet(tf, "Course/grade APIs", size=13)
add_bullet(tf, "Server-side rendering", size=13)

add_card(slide, 5.1, 3.2, 3, 2.8, RGBColor(0x15, 0x25, 0x10))
tf = add_text(slide, 5.3, 3.3, 2.6, 0.4, "LiveKit SFU :7880", size=16, color=GREEN, bold=True)
tf = add_text(slide, 5.3, 3.8, 2.6, 2, "", size=13, color=WHITE)
add_bullet(tf, "Video relay (SFU)", size=13)
add_bullet(tf, "Audio relay", size=13)
add_bullet(tf, "Screen share", size=13)
add_bullet(tf, "ICE/DTLS/TURN", size=13)
add_bullet(tf, "Simulcast layers", size=13)

add_card(slide, 8.8, 3.2, 3, 2.8)
tf = add_text(slide, 9.0, 3.3, 2.6, 0.4, "Go Server :4000", size=16, color=GOLD, bold=True)
tf = add_text(slide, 9.0, 3.8, 2.6, 2, "", size=13, color=WHITE)
add_bullet(tf, "WebSocket hub", size=13)
add_bullet(tf, "JWT token signing", size=13)
add_bullet(tf, "Breakout tokens", size=13)
add_bullet(tf, "Rate limiting", size=13)
add_bullet(tf, "Room management", size=13)

# Design principle
add_card(slide, 1.5, 6.3, 10.3, 0.9, RGBColor(0x1A, 0x15, 0x05))
add_text(slide, 1.8, 6.4, 9.7, 0.7, "Design Principle: LiveKit handles media (video/audio). Go WebSocket handles everything else (chat, whiteboard, polls, quiz). Clean separation = independent scaling.", size=14, color=GOLD)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: Features Beyond Video Conferencing
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0.8, 0.4, 11, 0.8, "9 Features That Zoom & Teams Don't Have", size=32, color=GOLD, bold=True)

features = [
    ("🖐 Hand Queue",       "Ordered queue with position tracking.\nTeacher calls next in order.", GREEN),
    ("🗳 Live Poll",         "Real-time voting with instant\nresult visualization.", BLUE),
    ("🧠 Quiz Mode",        "Timed in-class quiz with auto-scoring.\nFeeds into gradebook.", BLUE),
    ("👁 Focus Check",      "Anonymous comprehension pulse:\nClear / Confused / Too Fast.", GREEN),
    ("⏱ Pomodoro Timer",   "Class-wide 25/5 study timer\nsynced via WebSocket.", MUTED),
    ("🎯 Reactions",         "Floating educational emoji\nreactions over video grid.", MUTED),
    ("📝 Class Notes",       "Shared notepad. Teacher pins\nkey points. Students export.", GOLD),
    ("🚪 Breakout Rooms",   "Auto-group into N rooms.\nLiveKit sub-room switching.", GREEN),
    ("🎨 Whiteboard",       "Real-time collaborative canvas.\n1920×1080 fixed coordinates.", GOLD),
]

for i, (title, desc, color) in enumerate(features):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 1.4 + row * 2.0
    add_card(slide, x, y, 3.8, 1.7)
    add_text(slide, x + 0.2, y + 0.15, 3.4, 0.4, title, size=16, color=color, bold=True)
    add_text(slide, x + 0.2, y + 0.6, 3.4, 1, desc, size=12, color=WHITE)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: Technical Challenges
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0.8, 0.4, 11, 0.8, "Technical Challenges Overcome", size=32, color=GOLD, bold=True)

challenges = [
    ("WSL2 ICE Failure", "UDP ephemeral ports not forwarded through\nWSL2→Windows. ICE candidates unreachable.\nFix: force_tcp + built-in TURN server.", RED),
    ("React StrictMode", "Double-mount killed LiveKit WebSocket.\nSecond mount's cleanup closed first's connection.\nFix: Disabled StrictMode for classroom.", BLUE),
    ("Protocol Mismatch", "livekit-client v2.18.1 (protocol 16) vs\nserver v1.7.2 (protocol 13). JOIN_FAILURE.\nFix: Upgraded server to v1.10.1.", BLUE),
    ("Whiteboard Fragmentation", "Server rate limit (10 msg/s) silently dropped\n80% of stroke messages. Remote saw dots.\nFix: Client batching + increased limit to 30/s.", RED),
    ("Camera Persisting", "navigate() fired before React unmounted\nLiveKitRoom. MediaStreamTrack never stopped.\nFix: useEffect cleanup stops all tracks.", GOLD),
    ("Breakout Non-Functional", "Go server had event types but no handler.\nbreakout_join was just broadcast—no token.\nFix: Added JWT generation in WebSocket hub.", GOLD),
]

for i, (title, desc, color) in enumerate(challenges):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 1.4 + row * 2.9
    add_card(slide, x, y, 3.8, 2.5)
    add_text(slide, x + 0.2, y + 0.15, 3.4, 0.4, title, size=16, color=color, bold=True)
    add_text(slide, x + 0.2, y + 0.6, 3.4, 1.8, desc, size=12, color=WHITE)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: Debugging Deep Dive — The ICE Story
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0.8, 0.4, 11, 0.8, "Debugging Deep Dive: The ICE Candidate Mystery", size=32, color=GOLD, bold=True)
add_text(slide, 0.8, 1.1, 12, 0.4, "How I diagnosed the hardest bug — from 'leave-reconnect disconnected' to root cause in 6 steps", size=16, color=MUTED)

steps = [
    ("1. Browser Error (Unhelpful)", "leave-reconnect disconnected\ncould not establish pc connection\nNo useful information about WHY."),
    ("2. Server Log Analysis", "participant closing reason=SIGNAL_SOURCE_CLOSE\nsessionDuration: \"0s\"\nThe signal dies instantly — 0 second session!"),
    ("3. Transport Failure Found", "sent signal response: leave\nreason: CONNECTION_TIMEOUT\nclosing: TRANSPORT_FAILURE"),
    ("4. ICE Candidate Extraction", "Server: candidate udp 127.0.0.1:59847\nBrowser: candidate udp 172.18.80.1:53159\nNeither can reach the other!"),
    ("5. Root Cause: WSL2 Networking", "UDP ephemeral ports NOT forwarded.\n127.0.0.1 in browser = Windows, not WSL2.\n172.18.80.1 = Hyper-V adapter (unreachable)."),
    ("6. Fix: TCP-only + TURN", "force_tcp: true → stable port 7881\nturn.enabled: true → relay fallback\nConnection established, 0 → stable session."),
]

for i, (title, desc) in enumerate(steps):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 1.7 + row * 2.8
    c = GREEN if i == 5 else (RED if i < 2 else GOLD)
    add_card(slide, x, y, 3.8, 2.4)
    add_text(slide, x + 0.2, y + 0.15, 3.4, 0.4, title, size=14, color=c, bold=True)
    add_text(slide, x + 0.2, y + 0.55, 3.4, 1.7, desc, size=12, color=WHITE)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8: Performance Comparison
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0.8, 0.4, 11, 0.8, "Before & After: Performance Impact", size=32, color=GOLD, bold=True)

# Table
rows = [
    ("Metric",                   "P2P (Before)",        "LiveKit SFU (After)"),
    ("Max Participants",         "4-5",                 "100+"),
    ("Upload per User",          "N-1 streams",         "1 stream"),
    ("CPU per User",             "N-1 encodes",         "1 encode"),
    ("Screen Share Quality",     "720p",                "1080p @ 3Mbps"),
    ("Server Recording",         "Impossible",          "Built-in"),
    ("Late Join",                "Manual state sync",   "Automatic"),
    ("Breakout Rooms",           "Not possible",        "LiveKit sub-rooms"),
    ("Connection Reliability",   "Cascade failures",    "Server-managed"),
]

for i, (metric, before, after) in enumerate(rows):
    y = 1.5 + i * 0.58
    is_header = i == 0
    bg_c = RGBColor(0x1A, 0x2A, 0x1A) if is_header else (BG_CARD if i % 2 == 0 else BG_DARK)

    if is_header:
        add_card(slide, 0.8, y, 11.7, 0.5, RGBColor(0x1A, 0x2A, 0x1A))
    elif i % 2 == 0:
        add_card(slide, 0.8, y, 11.7, 0.5, BG_CARD)

    mc = GOLD if is_header else WHITE
    bc = GOLD if is_header else RED
    ac = GOLD if is_header else GREEN

    add_text(slide, 1.0, y + 0.05, 4, 0.4, metric, size=14, color=mc, bold=is_header)
    add_text(slide, 5.0, y + 0.05, 3.5, 0.4, before, size=14, color=bc, bold=is_header, align=PP_ALIGN.CENTER)
    add_text(slide, 8.5, y + 0.05, 4, 0.4, after, size=14, color=ac, bold=is_header, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9: Tech Stack & Skills Demonstrated
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0.8, 0.4, 11, 0.8, "Skills Demonstrated", size=32, color=GOLD, bold=True)

skills = [
    ("Frontend", ["React 18 (hooks, context, refs, error boundaries)",
                  "LiveKit components-react SDK integration",
                  "Real-time canvas drawing (HTML5 Canvas API)",
                  "WebSocket state management with monotonic sequencing",
                  "Responsive CSS Grid layout with dynamic pinning"], BLUE),
    ("Backend", ["Go (Gin framework) — WebSocket hub, JWT signing",
                 "Node.js/Express — auth, MongoDB, session management",
                 "Token-bucket rate limiting (custom implementation)",
                 "Gorilla WebSocket — per-client read/write pumps",
                 "Event-driven architecture (typed JSON envelopes)"], GREEN),
    ("Infrastructure", ["LiveKit SFU configuration & deployment",
                        "WebRTC fundamentals (ICE, DTLS, TURN, STUN)",
                        "Network debugging (WSL2 port forwarding, NAT)",
                        "JWT authentication (HMAC-SHA256, LiveKit claims)",
                        "Multi-service orchestration (3 servers + DB)"], GOLD),
]

for i, (category, items, color) in enumerate(skills):
    x = 0.8 + i * 4.1
    add_card(slide, x, 1.3, 3.8, 5.5)
    add_text(slide, x + 0.2, 1.4, 3.4, 0.5, category, size=20, color=color, bold=True)
    tf = add_text(slide, x + 0.2, 1.9, 3.4, 4.5, "", size=13, color=WHITE)
    for item in items:
        add_bullet(tf, "• " + item, size=12, color=WHITE)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10: Key Takeaways
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = GOLD; bar.line.fill.background()

add_text(slide, 0.8, 1.0, 11, 0.8, "Key Takeaways", size=36, color=GOLD, bold=True)

takeaways = [
    "Migrated from O(n²) peer-to-peer to O(n) SFU — enabling 100+ participant classrooms",
    "Built 9 education-specific features that Zoom/Teams don't offer",
    "Debugged complex WebRTC networking (ICE candidates, WSL2 port forwarding, protocol versions)",
    "Designed dual-transport architecture: LiveKit for media, Go WebSocket for features",
    "Solved real-time synchronization: whiteboard batching, rate limiting, coordinate systems",
    "Full-stack: React + Go + Node.js + LiveKit + MongoDB + WebSocket",
]

tf = add_text(slide, 1.5, 2.2, 10, 5, "", size=18, color=WHITE)
for t in takeaways:
    p = add_bullet(tf, t, size=18, color=WHITE)
    p.space_before = Pt(16)

add_text(slide, 1, 6.5, 11, 0.5, "github.com/cw2636/SLGS", size=16, color=MUTED, align=PP_ALIGN.CENTER)

bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.42), prs.slide_width, Inches(0.08))
bar2.fill.solid(); bar2.fill.fore_color.rgb = GOLD; bar2.line.fill.background()


# ═══════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════
output = "/home/wilson/SLGS/docs/SLGS_EduMeet_Technical_Presentation.pptx"
prs.save(output)
print(f"Saved: {output}")
